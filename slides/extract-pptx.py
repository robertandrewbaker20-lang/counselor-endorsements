#!/usr/bin/env python3
"""Extract slide text, notes, and theme colors from the Scouts AI PPTX."""

from __future__ import annotations

import json
import re
import sys
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

DEFAULT_PPTX = Path.home() / "Downloads" / "Scouts AI presentation 2026.pptx"
OUTPUT = Path(__file__).resolve().parent / "extracted-content.json"


def extract_text(shape) -> list[str]:
    texts: list[str] = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            t = "".join(run.text for run in para.runs).strip()
            if t:
                texts.append(t)
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            texts.extend(extract_text(child))
    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        for row in shape.table.rows:
            for cell in row.cells:
                t = cell.text.strip()
                if t:
                    texts.append(t)
    return texts


def theme_colors(pptx_path: Path) -> list[str]:
    colors: set[str] = set()
    with zipfile.ZipFile(pptx_path) as zf:
        for name in zf.namelist():
            if "theme" in name and name.endswith(".xml"):
                content = zf.read(name).decode("utf-8", errors="ignore")
                colors.update(re.findall(r'<a:srgbClr val="([0-9A-Fa-f]{6})"', content))
    return sorted(colors)


def main() -> int:
    pptx_path = Path(sys.argv[1]) if len(sys.argv) > 1 else DEFAULT_PPTX
    if not pptx_path.exists():
        print(f"File not found: {pptx_path}", file=sys.stderr)
        return 1

    prs = Presentation(str(pptx_path))
    slides_data = []

    for i, slide in enumerate(prs.slides, 1):
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        all_text: list[str] = []
        has_image = False
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                has_image = True
            all_text.extend(extract_text(shape))

        slides_data.append(
            {
                "slide_number": i,
                "layout_name": slide.slide_layout.name if slide.slide_layout else None,
                "all_text": all_text,
                "has_image": has_image,
                "speaker_notes": notes,
            }
        )

    output = {
        "source": str(pptx_path),
        "slide_count": len(slides_data),
        "theme_colors": theme_colors(pptx_path),
        "slides": slides_data,
    }

    OUTPUT.write_text(json.dumps(output, indent=2))
    print(f"Wrote {OUTPUT} ({len(slides_data)} slides)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
